import pprint
from datetime import date
# Import the byte stream handler.
from io import BytesIO
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE, MSO_UNDERLINE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR


base_url = "https://api.trydatabook.com"


def get_token(base):
    auth_url = base + "/auth/local"
    res = requests.post(auth_url, json={"email": "nikhilts19@gmail.com", "password": "Mypassword1"})
    return res.json()["token"]


token = get_token(base_url)


def get_details_page2(base, company_id, token):
    # try:

    # token = get_token()
    response = requests.get(base + "/api/companies/" + company_id, headers={'Authorization': "Bearer " + token})
    print(response.status_code)
    data = response.json()
    pprint.pprint(data)
    return data['description'], data['foundedYear'], data['type'], data['latestRevenue']['valueUSD'], \
           data['latestRevenue']['currency'], \
           data['headquarters']['address']
    # print(data["_id"])

    # except HTTPStatus.OK == 401:
    #     token = get_token()


# def iter_cells(table):
#     for row in table.rows:
#         for cell in row.cells:
#             yield cell


def create_ppt_page2():
    # prs = Presentation("test.pptx")
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    # title_only_slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(title_only_slide_layout)
    # shapes = slide.shapes
    #
    # shapes.title.text = 'Adding a Table'
    #
    # rows = 15
    # cols = 6
    # left = top = Inches(0.5)
    # right = Inches(0.5)
    # width = Inches(5.0)
    # height = Inches(0.2)
    #
    # table = shapes.add_table(rows, cols, left, top, width, height).table
    #
    # # set column widths
    # table.columns[0].width = Inches(2.0)
    # table.columns[1].width = Inches(4.0)
    #
    # # write column headings
    # table.cell(0, 0).text = 'Category of KRI '
    # table.cell(0, 1).text = 'Defintion'
    # table.cell(0, 2).text = 'Measure'
    # table.cell(0, 3).text = 'Previous Score'
    # table.cell(0, 4).text = 'Current Score'
    # table.cell(0, 5).text = 'Risk Trend'
    #
    # # write body cells
    # # table.cell(2, 0).text = 'Inventory'
    # for pp in table.cell(1, 0).text_frame.paragraphs:
    #     for run in pp.runs:
    #         run.font.bold = True
    #
    # table.cell(1, 1).text = 'Qux'
    #
    # for cell in iter_cells(table):
    #     for paragraph in cell.text_frame.paragraphs:
    #         for run in paragraph.runs:
    #             run.font.size = Pt(12)
    #
    # table.cell(1, 0).text = 'Random'
    # table.cell(1, 1).text = 'Random'
    # table.cell(1, 2).text = 'Random'
    # table.cell(1, 3).text = 'Random'
    # table.cell(1, 4).text = 'Random'
    #
    # # table.cell(2,0).text = 'Inventory'
    # table.cell(5, 0).text = 'random'
    # table.cell(6, 0).text = 'Random'
    # table.cell(7, 0).text = 'Random'
    #
    # col_cells = []
    #
    # col_cells.append(table.cell(1, 0))

    description, founded_year, type_of_company, latest_revenue_value, latest_revenue_currency, address = get_details_page2(
        base_url,
        str(company_id),
        token)
    prs = Presentation("test.pptx")

    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Company introduction'

    tf = body_shape.text_frame
    # tf.autofit_text()
    # tf.fit_text()
    tf.text = 'Description of the company'
    tf.word_wrap = True

    p = tf.add_paragraph()
    # description_split = description[:550].split('.')
    # p.text = " ".join(description_split[:len(description_split)-1])
    p.text = description
    p.alignment = PP_ALIGN.JUSTIFY

    p.word_wrap = True
    p.level = 1

    font = p.font
    # font.name = 'Calibri'
    font.size = Pt(14)
    # font.bold = True
    font.italic = None  # cause value to be inherited from theme
    # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    p.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Page 3 generation of details of the company
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Company introduction continued'

    tf = body_shape.text_frame
    tf.text = 'Details of the company'
    l1 = [('Year found', founded_year), ('Company type', type_of_company),
          ('Latest_revenue in millions', latest_revenue_value), ('Revenue currency', latest_revenue_currency),
          ('Address', address)]
    for t, val in l1:
        p = tf.add_paragraph()
        p.text = t + " : " + str(val)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    # p = tf.add_paragraph()
    # p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    # p.level = 2

    # # adding bg image
    # bg_img = requests.get('https://cdn.hipwallpaper.com/i/95/32/hMtTPx.jpg')
    # bg_image_data = BytesIO(bg_img.content)
    # left = top = 0
    # bg = slide.shapes.add_picture(bg_image_data, left - 0.1 * prs.slide_width, top, height=prs.slide_height)

    # title.text = "Company introduction"
    # subtitle.text = description + "\nGenerated on {:%m-%d-%Y}".format(date.today())

    prs.save('test.pptx')

    print("PPT page 2 saved to the folder!!!!")


company_id = "58153a9ac12a6e000f7c2bea"
# print(get_details(base_url, company_id, token))
create_ppt_page2()
