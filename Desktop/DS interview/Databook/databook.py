import pprint
from datetime import date
from io import BytesIO

import requests
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

import config
import create_graph

# use caps for global variables
BASE_URL = "https://api.trydatabook.com"

company_id = "584d79ba083592000413cf7a"  # "58153a9ac12a6e000f7c2bea"
metric_id = "57fe09db3e767b000f54a659"


def get_token(base):
    auth_url = base + "/auth/local"
    res = requests.post(auth_url, json={"email": config.EMAIL_ADDRESS, "password": config.PASSWORD})
    return res.json()["token"]


# token = get_token(BASE_URL)
token = 'eyJhbGciOiJSUzI1NiIsInR5cCI6IkpXVCIsImtpZCI6IlFVSXhSVFk1TXpCR09URTVPVVpFTURJNU1ERTFNRGMwTnpNeE1FSkdOamxDUlVSQ1JFUXpRUSJ9.eyJpc3MiOiJodHRwczovL2RhdGFib29rLmF1dGgwLmNvbS8iLCJzdWIiOiJhdXRoMHxkYXRhYm9va3w1ZTNjY2ZlOWI0ZDEzYTAwNDlmNzE2ODgiLCJhdWQiOiJodHRwczovL2FwaS50cnlkYXRhYm9vay5jb20iLCJpYXQiOjE1ODE5MTU0MTMsImV4cCI6MTU4MjAwMTgxMywiYXpwIjoic2lSVllIaGVWc3ExMEkwdDQzNnNwczg1Sk1KMnphakoiLCJzY29wZSI6InN1cGVyQWRtaW4gYWRtaW4iLCJndHkiOiJwYXNzd29yZCJ9.GX3RIbF09Aqn7PbnbsXIOXsEYvs7MYdIbEyWCXbBZjp10-KiLjYVT1yRJTyi94Y_9L38LSnEGBjnS3hpX6LRnKLtTY08FQk46O0UFUOrMbjkptfG65uCXaUTvNUJf_ZezhcGV1IT2S7Ij6qtGf9g3dXu8eFlQWnrgo9wmnsh_yFiS1K4sQpewrlgP27-190SM3O1bYAFmvkQUFoNWPsAokOZdXNLYYGBz5CossOM0DzGEpfWy7uNE_d-vAZQS-9RFqtPoK_qPhH4ov75o0rxiT6WL_FNrSB0eaVcCEH-2q1b0jyi8GvPizZEJnjpswGYrrIkpUfEdSVSp9y269tN6Q '


# print(token)


def get_details(base, company_id, token):
    # do not use generic variable names
    response = requests.get(base + "/api/companies/" + company_id, headers={'Authorization': "Bearer " + token})
    print(response.status_code)
    data = response.json()
    pprint.pprint(data)
    return data


def get_details_graph(base, company_id, token, metric_id):
    # do not use generic variable names
    response = requests.get(base + "/api/companies/" + company_id + "/metrics/" + metric_id,
                            headers={'Authorization': "Bearer " + token})
    print(response.status_code)
    data = response.json()
    pprint.pprint(data)
    return data


chart_data = get_details_graph(BASE_URL, company_id, token, metric_id)
company_metadata = get_details(BASE_URL, company_id, token)


def create_ppt_page_one(data):
    name, website, logo_url, latest_revenue_growth = data['name'], data['website'], data['logoUrl'], data[
        'latestRevenueGrowth']
    prs = Presentation("test.pptx")
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = name
    subtitle.text = "\nWebsite : " + website + "\ndata upto " + str(latest_revenue_growth['year']) + " quarter " + str(
        latest_revenue_growth['quarter']) + "\nGenerated on {:%m-%d-%Y}".format(date.today())

    # Adding the logo image to the first page
    r = requests.get(logo_url)
    image_data = BytesIO(r.content)
    left = top = Inches(1)
    pic = slide.shapes.add_picture(image_data, left, top)

    prs.save('test.pptx')

    print("First page is created!!!!\n")


def create_ppt_page_two(data):
    description, founded_year, type_of_company, latest_revenue_value, latest_revenue_currency, address = data[
                                                                                                             'companyDescription'], \
                                                                                                         data[
                                                                                                             'foundedYear'], \
                                                                                                         data['type'], \
                                                                                                         data[
                                                                                                             'latestRevenue'][
                                                                                                             'valueUSD'], \
                                                                                                         data[
                                                                                                             'latestRevenue'][
                                                                                                             'currency'], \
                                                                                                         data[
                                                                                                             'headquarters'][
                                                                                                             'address']

    prs = Presentation("test.pptx")

    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Company introduction'

    tf = body_shape.text_frame
    tf.text = 'Description of the company'
    tf.word_wrap = True

    p = tf.add_paragraph()
    # description_split = description[:550].split('.')
    # p.text = " ".join(description_split[:len(description_split)-1])
    p.text = description
    p.alignment = PP_ALIGN.JUSTIFY
    p.word_wrap = True
    p.level = 1

    # font = p.font
    # # font.name = 'Calibri'
    # font.size = Pt(14)
    # # font.bold = True
    # font.italic = None  # cause value to be inherited from theme
    # # font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # p.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    # Page 3 generation of details of the company
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Company introduction continued'

    tf = body_shape.text_frame
    tf.text = 'Details of the company'
    l1 = [('Year found', founded_year), ('Company type', type_of_company),
          ('Latest_revenue in millions', latest_revenue_value), ('Revenue currency', latest_revenue_currency),
          ('Address', address), ('CEO', data['ceo']), ('Country', data['country'])]
    for t, val in l1:
        p = tf.add_paragraph()
        p.text = t + " : " + str(val)
        p.alignment = PP_ALIGN.LEFT
        p.level = 1

    prs.save('test.pptx')

    print("PPT page 2 saved to the folder!!!!")


yearly_data, quarterly_data, company_name = create_graph.get_graph_data(chart_data, company_id)
create_graph.generate_graph(yearly_data, quarterly_data, company_name)


def create_yearly_revenue_graphs(chart_data, chart):
    prs = Presentation("test.pptx")
    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Yearly revenue of the company"
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(chart)
    pic.crop_top = 0
    pic.crop_left = 0
    pic.crop_bottom = 0
    pic.crop_right = 0
    print("======> picture 1 inserted")

    subtitle = slide.placeholders[2]
    subtitle.text = chart_data['description']
    subtitle.alignment = PP_ALIGN.JUSTIFY

    prs.save('test.pptx')

    print("Chart PPT saved to the folder!!!!")


create_ppt_page_one(company_metadata)
create_ppt_page_two(company_metadata)
for chart in ['yearly_chart.jpg', '2019_quarterly_chart.jpg']:
    create_yearly_revenue_graphs(chart_data, chart)
