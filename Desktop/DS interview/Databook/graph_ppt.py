import json

import matplotlib
import pandas as pd
import matplotlib.pyplot as plt
# %matplotlib inline
import pprint
from datetime import date
from io import BytesIO
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_UNDERLINE, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

# use caps for global variables
# use Underscore(_) for namings and be consistent
BASE_URL = "https://api.trydatabook.com"
company_id = "58153a9ac12a6e000f7c2bea"
metric_id = "57fe09db3e767b000f54a659"


def get_token(base):
    auth_url = base + "/auth/local"
    res = requests.post(auth_url, json={"email": "nikhilts19@gmail.com", "password": "Mypassword1"})
    return res.json()["token"]


token = get_token(BASE_URL)


def get_details_page1(base, company_id, token, metric_id):
    # do not use generic variable names
    response = requests.get(base + "/api/companies/" + company_id + "/metrics/" + metric_id,
                            headers={'Authorization': "Bearer " + token})
    print(response.status_code)
    data = response.json()
    pprint.pprint(data)
    return data


chart_data = get_details_page1(BASE_URL, company_id, token, metric_id)
# with open('chart_data.json', 'w') as json_file:
#     json.dump(chart_data, json_file)
#     print("Chart Data saved to chart_data JSON file!!!!")


def create_ppt_page_graphs(chart_data, chart):
    prs = Presentation("test.pptx")
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a chart slide first
    # Create the summary graph
    graph_slide_layout = prs.slide_layouts[8]
    slide = prs.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text = "Revenue by company"
    placeholder = slide.placeholders[1]
    pic = placeholder.insert_picture(chart)
    subtitle = slide.placeholders[2]
    subtitle.text = chart_data['description']
    subtitle.alignment = PP_ALIGN.JUSTIFY

    prs.save('test.pptx')

    print("Chart PPT saved to the folder!!!!")


# get_details_page1(BASE_URL, company_id, token, metric_id)
create_ppt_page_graphs(chart_data, 'data.jpg')